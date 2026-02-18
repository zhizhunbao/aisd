"""
Extract text AND images from PPTX slides for midterm review.
Processes all 5 weeks of MV course slides.
Images are saved per-slide and referenced in the markdown.
"""
import os
import io
import hashlib
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SLIDES_DIR = Path(r"d:\BaiduSyncdisk\workspace\python_workspace_2025\aisd\courses\mv\slides")
NOTES_DIR = Path(r"d:\BaiduSyncdisk\workspace\python_workspace_2025\aisd\courses\mv\notes")

# Week mapping: filename -> (week_number, topic_slug)
WEEKS = {
    "Week 1 - Introduction to Machine Vision1.pptx": (1, "week1_intro_machine_vision"),
    "Week 2 - Fundamentals of Image Processing1.pptx": (2, "week2_image_processing"),
    "Week. 3-Object_Feature Detection and Description.pptx": (3, "week3_feature_detection"),
    "Week 4 - Introduction to Convolutional Neural Networks (CNNs)1.pptx": (4, "week4_cnn"),
    "Week5_ Deep Learning for Image Classification1.pptx": (5, "week5_deep_learning"),
}

# Skip tiny images (logos, bullets, decorations) by minimum size in bytes
MIN_IMAGE_BYTES = 2000


def extract_images_from_slide(slide, slide_idx, img_dir, topic):
    """Extract all meaningful images from a slide, return markdown image references."""
    image_refs = []
    img_count = 0
    seen_hashes = set()

    for shape in slide.shapes:
        # Direct picture shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            img_bytes = image.blob
            
            # Skip tiny images (icons, bullets)
            if len(img_bytes) < MIN_IMAGE_BYTES:
                continue
            
            # Deduplicate by content hash
            img_hash = hashlib.md5(img_bytes).hexdigest()[:8]
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)
            
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            
            img_count += 1
            img_filename = f"slide{slide_idx:02d}_img{img_count}.{ext}"
            img_path = img_dir / img_filename
            img_path.write_bytes(img_bytes)
            
            # Use alt text if available, otherwise generic label
            alt_text = shape.name or f"Slide {slide_idx} Image {img_count}"
            rel_path = f"{img_dir.name}/{img_filename}"
            image_refs.append(f"![{alt_text}]({rel_path})")
        
        # Group shapes may contain pictures
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = child.image
                    img_bytes = image.blob
                    if len(img_bytes) < MIN_IMAGE_BYTES:
                        continue
                    img_hash = hashlib.md5(img_bytes).hexdigest()[:8]
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_count += 1
                    img_filename = f"slide{slide_idx:02d}_img{img_count}.{ext}"
                    img_path = img_dir / img_filename
                    img_path.write_bytes(img_bytes)
                    alt_text = child.name or f"Slide {slide_idx} Image {img_count}"
                    rel_path = f"{img_dir.name}/{img_filename}"
                    image_refs.append(f"![{alt_text}]({rel_path})")

    return image_refs


def extract_slide_text(pptx_path: Path, week_num: int, topic: str) -> tuple:
    """Extract all text and images from a PPTX file into structured markdown."""
    prs = Presentation(str(pptx_path))
    
    # Create image directory
    img_dir = NOTES_DIR / f"{topic}_slides_images"
    img_dir.mkdir(parents=True, exist_ok=True)
    
    lines = []
    lines.append(f"# Week {week_num}: {pptx_path.stem}")
    lines.append("")
    lines.append(f"> Source: `{pptx_path.name}`")
    lines.append(f"> Total slides: {len(prs.slides)}")
    lines.append("")
    lines.append("---")
    lines.append("")
    
    total_images = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_idx}")
        lines.append("")

        # Extract text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if para.level == 0 and shape == slide.shapes[0]:
                            lines.append(f"### {text}")
                        elif para.level > 0:
                            indent = "  " * (para.level - 1)
                            lines.append(f"{indent}- {text}")
                        else:
                            lines.append(text)
                lines.append("")
            
            elif shape.has_table:
                table = shape.table
                header = " | ".join(cell.text.strip() for cell in table.rows[0].cells)
                lines.append(f"| {header} |")
                lines.append("|" + "---|" * len(table.rows[0].cells))
                for row in list(table.rows)[1:]:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    lines.append(f"| {row_text} |")
                lines.append("")

        # Extract images
        image_refs = extract_images_from_slide(slide, slide_idx, img_dir, topic)
        if image_refs:
            total_images += len(image_refs)
            for ref in image_refs:
                lines.append(ref)
                lines.append("")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append(f"> **Speaker Notes:** {notes_text}")
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines), total_images, len(prs.slides)


def main():
    NOTES_DIR.mkdir(parents=True, exist_ok=True)
    
    print("=" * 60)
    print("  MV Course Slides Extraction (Text + Images)")
    print("=" * 60)
    print()
    
    for filename, (week_num, topic) in WEEKS.items():
        pptx_path = SLIDES_DIR / filename
        if not pptx_path.exists():
            print(f"❌ Not found: {filename}")
            continue
        
        output_path = NOTES_DIR / f"{topic}_slides.md"
        
        print(f"[{week_num}/5] {filename}")
        print(f"       Size: {pptx_path.stat().st_size / 1024 / 1024:.1f} MB")
        
        content, img_count, slide_count = extract_slide_text(pptx_path, week_num, topic)
        output_path.write_text(content, encoding="utf-8")
        
        print(f"       → {output_path.name} ({len(content):,} chars)")
        print(f"       → {topic}_slides_images/ ({img_count} images from {slide_count} slides)")
        print()
    
    # Summary
    print("=" * 60)
    total_imgs = sum(
        len(list((NOTES_DIR / f"{t}_slides_images").glob("*")))
        for _, (_, t) in WEEKS.items()
        if (NOTES_DIR / f"{t}_slides_images").exists()
    )
    print(f"  ✅ Done! Total images extracted: {total_imgs}")
    print("=" * 60)


if __name__ == "__main__":
    main()
