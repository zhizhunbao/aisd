import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn

prs = Presentation(r'c:\Users\40270\Desktop\workspace\aisd\courses\mv\project\CST8508_FinalProject_Presentation_HyeRanYoo_PengWang_final.pptx')

for si in [14]:  # slide 15 (Summary)
    slide = prs.slides[si]
    print(f'=== Slide {si+1}: ALL shapes ===')
    for s in slide.shapes:
        sp = s._element
        prstGeom = sp.find(f'.//{qn("a:prstGeom")}')
        geom_name = prstGeom.get('prst') if prstGeom is not None else 'none'
        solid = sp.find(f'.//{qn("a:solidFill")}')
        fill_color = ""
        if solid is not None:
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is not None:
                fill_color = f'#{srgb.get("val")}'
        text = s.text_frame.text[:60] if s.has_text_frame else ""
        print(f'  {s.name} | geom={geom_name} | fill={fill_color} | L={s.left} T={s.top} W={s.width} H={s.height} | "{text}"')
