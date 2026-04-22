"""
Update Final Project Presentation PPTX — V3
============================================
Changes:
1. Slide 4  – pin exact dependency versions
2. Slide 14 – redesign "Challenges & Lessons Learned" (two-column + reasons)
3. NEW slide – "Tech Stack" version table (inserted after slide 14)
4. Slide 15→16 – replace Summary with "Evaluation" (strengths + limitations)
5. Q&A remains last

Final order: ... | 12. Challenges & LL | 13. Tech Stack | 14. Evaluation | 15. Q&A
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

SRC = r"c:\Users\40270\Desktop\workspace\aisd\courses\mv\project\CST8508_FinalProject_Presentation_HyeRanYoo_PengWang_final.pptx"
DST = SRC.replace(".pptx", "_updated.pptx")

prs = Presentation(SRC)

# ─── Design constants ────────────────────────────────────────────
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
GREEN      = RGBColor(0x00, 0xC9, 0x7B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
DARK_BG    = RGBColor(0x2D, 0x2D, 0x44)
BLUE_ACC   = RGBColor(0x00, 0x96, 0xD6)
OVAL_SIZE  = 127000

# ─── Helpers ─────────────────────────────────────────────────────
def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def set_text_keep_format(shape, new_text):
    if shape and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
                return
        shape.text_frame.paragraphs[0].text = new_text

def remove_shapes_by_names(slide, names):
    sp_tree = slide.shapes._spTree
    to_remove = [s._element for s in slide.shapes if s.name in names]
    for el in to_remove:
        sp_tree.remove(el)

def add_oval(slide, left, top, color, size=OVAL_SIZE):
    shape = slide.shapes.add_shape(9, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text,
                font_size=14, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def add_bullet_item(slide, left_oval, left_text, top, width,
                    main_text, sub_text, oval_color=ORANGE,
                    main_size=13, sub_size=10):
    add_oval(slide, left_oval, top + Pt(5), oval_color)
    main_h = Pt(22)
    add_textbox(slide, left_text, top, width, main_h,
                main_text, font_size=main_size, color=WHITE)
    if sub_text:
        sub_top = top + main_h + Pt(1)
        add_textbox(slide, left_text + Pt(8), sub_top, width - Pt(8), Pt(18),
                    sub_text, font_size=sub_size, color=LIGHT_GRAY)
    return main_h + (Pt(22) if sub_text else Pt(4))

def add_slide_background(slide):
    """Add the standard dark background + header + blue accent bar (matching all slides)."""
    # Slide-level solid fill (matches all existing slides: #1A1A2E)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    # Dark header rect
    rect1 = slide.shapes.add_shape(1, 0, 0, prs.slide_width, 1188720)
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = DARK_BG
    rect1.line.fill.background()
    # Blue accent bar
    rect2 = slide.shapes.add_shape(1, 0, 1188720, prs.slide_width, 45720)
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = BLUE_ACC
    rect2.line.fill.background()

def move_slide_to_index(prs, slide, new_index):
    """Move a slide to a specific position (0-indexed)."""
    # Access the presentation XML element
    pres_elm = prs.part._element
    sldIdLst = pres_elm.find(qn('p:sldIdLst'))
    # Find rId for this slide
    slide_rId = None
    for rel_key, rel in prs.part.rels.items():
        if rel.target_part == slide.part:
            slide_rId = rel_key
            break
    if not slide_rId:
        print("  WARNING: Could not find slide rId")
        return
    # Find the sldId element
    target_el = None
    for sldId in sldIdLst:
        if sldId.get(qn('r:id')) == slide_rId:
            target_el = sldId
            break
    if target_el is None:
        print("  WARNING: Could not find sldId element")
        return
    sldIdLst.remove(target_el)
    all_ids = list(sldIdLst)
    if new_index >= len(all_ids):
        sldIdLst.append(target_el)
    else:
        all_ids[new_index].addprevious(target_el)


# ═══════════════════════════════════════════════════════════════════
# 1. SLIDE 4: Pin exact versions
# ═══════════════════════════════════════════════════════════════════
slide4 = prs.slides[3]
for name, text in {
    "TextBox 13": "3.11.15 (3.13 incompatible with MMCV)",
    "TextBox 15": "2.2.2+cpu / torchvision 0.17.2+cpu",
    "TextBox 19": "3.2.0 / mmengine 0.10.7",
    "TextBox 21": "0.27.0 ByteTrack tracker (CPU)",
    "TextBox 23": "4.13.0 Webcam capture & display",
    "TextBox 25": "1.26.4 Array operations",
}.items():
    shape = find_shape(slide4, name)
    if shape:
        set_text_keep_format(shape, text)
print("Slide 4: Versions updated")


# ═══════════════════════════════════════════════════════════════════
# 2. SLIDE 14: Challenges & Lessons Learned (two-column)
# ═══════════════════════════════════════════════════════════════════
slide14 = prs.slides[13]
set_text_keep_format(find_shape(slide14, "TextBox 3"), "12. Challenges & Lessons Learned")
set_text_keep_format(find_shape(slide14, "TextBox 4"), "What we faced & what we gained")

# Remove old content
remove_shapes_by_names(slide14, [
    "TextBox 17", "Oval 18", "TextBox 19", "Oval 20", "TextBox 21",
    "Oval 22", "TextBox 23", "Oval 24", "TextBox 25"
])

COL1_OVAL = 792480;   COL1_TEXT = 1066800;  COL1_W = 4800000
COL2_OVAL = 6200000;  COL2_TEXT = 6474320;  COL2_W = 4800000
TOP = 1500000

# Left: Challenges
add_textbox(slide14, COL1_OVAL, TOP, COL1_W, Pt(24),
            "Challenges", font_size=18, color=ORANGE, bold=True)
y = TOP + Pt(32)
for main, sub in [
    ("mmdetection 3.2.0 Windows install",
     "C++ extensions require exact compiler & version matching"),
    ("Python 3.13 breaks MMCV 2.2.0",
     "Pre-built wheels only exist for Python <=3.11"),
    ("Camera warm-up: dark first frames",
     "Auto-exposure needs 10-20 frames to stabilize"),
    ("Confidence threshold balancing",
     "0.3 = false positives, 0.7 = misses distant persons, 0.5 optimal"),
]:
    h = add_bullet_item(slide14, COL1_OVAL, COL1_TEXT, y, COL1_W,
                        main, sub, ORANGE, 13, 10)
    y += h + Pt(10)

# Right: Lessons Learned
add_textbox(slide14, COL2_OVAL, TOP, COL2_W, Pt(24),
            "Lessons Learned", font_size=18, color=GREEN, bold=True)
y = TOP + Pt(32)
for main, sub in [
    ("End-to-end CV pipeline thinking",
     "Each module's output feeds the next; one failure cascades"),
    ("Version matrix is critical",
     "PyTorch 2.2.2 - MMCV 2.2.0 - mmdet 3.2.0 must match exactly"),
    ("Config/code separation",
     "External JSON lets you tune thresholds without touching code"),
    ("Real-time debugging needs HUD",
     "Can't set breakpoints in video loop; visual overlay is essential"),
]:
    h = add_bullet_item(slide14, COL2_OVAL, COL2_TEXT, y, COL2_W,
                        main, sub, GREEN, 13, 10)
    y += h + Pt(10)
print("Slide 14: Challenges & Lessons Learned (two-column)")


# ═══════════════════════════════════════════════════════════════════
# 3. NEW SLIDE: Tech Stack version table (insert after slide 14)
# ═══════════════════════════════════════════════════════════════════
# Use blank layout
blank_layout = None
for layout in prs.slide_layouts:
    if layout.name == 'Blank':
        blank_layout = layout
        break
if blank_layout is None:
    blank_layout = prs.slide_layouts[6]  # fallback to blank

new_slide = prs.slides.add_slide(blank_layout)

# Background
add_slide_background(new_slide)

# Title
add_textbox(new_slide, 731520, 182880, 9144000, Pt(44),
            "13. Tech Stack", font_size=36, color=WHITE, bold=True)
add_textbox(new_slide, 731520, 685800, 9144000, Pt(24),
            "Verified package versions from project environment",
            font_size=16, color=LIGHT_GRAY)

# Presenter tag
add_textbox(new_slide, 7772400, 777240, 3657600, Pt(22),
            "Presenter: Peng Wang", font_size=12, color=GREEN)

# --- Create table ---
rows = 11  # header + 10 packages
cols = 3
table_left   = 914400
table_top    = 1450000
table_width  = 10363200
table_height = Pt(rows * 30)

table_shape = new_slide.shapes.add_table(rows, cols, table_left, table_top,
                                          table_width, table_height)
table = table_shape.table

# Column widths
table.columns[0].width = 3200000   # Package
table.columns[1].width = 2800000   # Version
table.columns[2].width = 4363200   # Role

# Table data
data = [
    ("Package",       "Version",         "Role in Project"),
    ("Python",        "3.11.15",         "Runtime (3.13 incompatible with MMCV C++ extensions)"),
    ("PyTorch",       "2.2.2+cpu",       "Deep learning framework for RTMDet-s inference"),
    ("torchvision",   "0.17.2+cpu",      "Image transforms & model utilities"),
    ("MMCV",          "2.2.0",           "OpenMMLab foundation library (pre-built CPU wheel)"),
    ("mmdetection",   "3.2.0",           "Object detection toolbox providing RTMDet-s model"),
    ("mmengine",      "0.10.7",          "OpenMMLab training/inference engine"),
    ("supervision",   "0.27.0",          "ByteTrack multi-object tracker implementation"),
    ("OpenCV",        "4.13.0",          "Webcam capture, frame display, image drawing"),
    ("NumPy",         "1.26.4",          "Array operations for bbox math & filtering"),
]

# Style the table
HEADER_BG = RGBColor(0x00, 0x96, 0xD6)  # blue accent
ROW_BG_1  = RGBColor(0x2A, 0x2A, 0x3E)
ROW_BG_2  = RGBColor(0x33, 0x33, 0x4D)

def set_cell(cell, text, font_size=12, bold=False, color=WHITE, bg_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Consolas" if not bold else None
    # Cell background
    if bg_color:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', f'{bg_color}')
    # Remove cell margins for compact look
    cell.margin_left = Pt(8)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        if row_idx == 0:
            set_cell(cell, text, font_size=13, bold=True, color=WHITE,
                     bg_color='0096D6')
        else:
            bg = '2A2A3E' if row_idx % 2 == 1 else '33334D'
            # Version column in green
            c = GREEN if col_idx == 1 else WHITE
            set_cell(cell, text, font_size=11, bold=(col_idx == 0),
                     color=c, bg_color=bg)

# Remove table borders for clean look
for row in table.rows:
    for cell in row.cells:
        tcPr = cell._tc.get_or_add_tcPr()
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = etree.SubElement(tcPr, qn(border_name))
            ln.set('w', '0')
            noFill = etree.SubElement(ln, qn('a:noFill'))

print("New slide: Tech Stack table created")

# --- Move new slide to position 14 (after Challenges, before old Summary/Eval) ---
# New slide was appended at end (index 16). Move to index 14.
move_slide_to_index(prs, new_slide, 14)
print("New slide: Moved to position 15 (after Challenges)")


# ═══════════════════════════════════════════════════════════════════
# 4. SLIDE 15→16: Replace Summary with Evaluation
# ═══════════════════════════════════════════════════════════════════
# After inserting new slide, old slide 15 (Summary) is now at index 15
slide_eval = prs.slides[15]
set_text_keep_format(find_shape(slide_eval, "TextBox 3"), "14. Evaluation")
set_text_keep_format(find_shape(slide_eval, "TextBox 4"), "System Performance & Assessment")

# Remove old Summary bullet items
old_names = [f"Oval {i}" for i in range(5, 20)] + [f"TextBox {i}" for i in range(6, 21)]
remove_shapes_by_names(slide_eval, old_names)

EVAL_TOP = 1500000

# Left: Strengths
add_textbox(slide_eval, COL1_OVAL, EVAL_TOP, COL1_W, Pt(24),
            "Strengths", font_size=18, color=GREEN, bold=True)
y = EVAL_TOP + Pt(32)
for main, sub in [
    ("RTMDet-s: 44.6% AP on COCO, CPU real-time",
     "Pre-trained model, no GPU or fine-tuning required"),
    ("ByteTrack (supervision 0.27.0): stable IDs",
     "Two-stage IoU association + 30-frame occlusion buffer"),
    ("5% dead zone prevents direction jitter",
     "Only fires when offset exceeds frame_size x 0.05"),
    ("Modular architecture",
     "Can swap RTMDet-s for YOLOv8 or ByteTrack for DeepSORT"),
    ("Dual-window: Full View + Person Focus",
     "Context + detail simultaneously, 640x480 + 480x480"),
]:
    h = add_bullet_item(slide_eval, COL1_OVAL, COL1_TEXT, y, COL1_W,
                        main, sub, GREEN, 13, 10)
    y += h + Pt(8)

# Right: Limitations
add_textbox(slide_eval, COL2_OVAL, EVAL_TOP, COL2_W, Pt(24),
            "Limitations", font_size=18, color=RED_ACCENT, bold=True)
y = EVAL_TOP + Pt(32)
for main, sub in [
    ("CPU-only: ~10-15 FPS",
     "GPU (CUDA) would push to 30+ FPS for production"),
    ("No depth estimation",
     "Single 2D camera can't determine real-world distance"),
    ("Largest bbox heuristic can fail",
     "Closer bystander takes over when target walks away"),
    ("No re-identification (ReID)",
     "Target leaving & returning gets assigned a new tracker ID"),
    ("No servo/motor integration",
     "Direction messages are visual only; no physical PTZ control"),
]:
    h = add_bullet_item(slide_eval, COL2_OVAL, COL2_TEXT, y, COL2_W,
                        main, sub, RED_ACCENT, 13, 10)
    y += h + Pt(8)

print("Slide 16: Evaluation (Strengths + Limitations)")


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
prs.save(DST)
print(f"\n✅ Saved: {DST}")
print("Slide order: ... | 12. Challenges & LL | 13. Tech Stack | 14. Evaluation | 15. Q&A")
